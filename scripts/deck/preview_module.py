"""Build ONE slides module in isolation and render its slides to PNG for visual checking.

Usage: python scripts/deck/preview_module.py deck.slides_answers /tmp/.../outdir
Renders with spire.presentation (evaluation watermark top-left is expected in previews only).
"""
import importlib
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
from pptx import Presentation  # noqa: E402
from deck import config as C  # noqa: E402
from deck import helpers as H  # noqa: E402


def main(modname, outdir):
    os.makedirs(outdir, exist_ok=True)
    mod = importlib.import_module(modname)
    prs = Presentation()
    prs.slide_width, prs.slide_height = C.SL_W, C.SL_H
    H.set_total(len(mod.SLIDES))
    for fn in mod.SLIDES:
        fn(prs)
    pptx = os.path.join(outdir, modname.split(".")[-1] + ".pptx")
    prs.save(pptx)
    print("saved", pptx)
    try:
        from spire.presentation import Presentation as SP
        p = SP()
        p.LoadFromFile(pptx)
        for i in range(p.Slides.Count):
            img = p.Slides[i].SaveAsImage()
            out = os.path.join(outdir, f"slide_{i+1:02d}.png")
            img.Save(out)
            img.Dispose()
            print("rendered", out)
        p.Dispose()
    except Exception as e:  # pragma: no cover
        print("render skipped:", e)
    # word / number budget report
    import re
    prs2 = Presentation(pptx)
    for i, s in enumerate(prs2.slides, 1):
        txt = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        words = len(txt.split())
        nums = re.findall(r"(?<![A-Za-z_])[-+]?\d+(?:\.\d+)?(?![A-Za-z_])", txt)
        notes = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        print(f"slide {i:2d}: {words:3d} words, {len(nums):2d} numeric tokens, notes {len(notes):4d} chars")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
