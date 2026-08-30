"""Build the 2026-08-31 supervisor-meeting deck (pptx) from the modular slide builders.

Usage:  python scripts/deck/build_deck_2026_08_31.py [--render DIR] [--skip-missing]
Output: notes/supervisor_meeting_2026_08_31.pptx  (+ copy figures/supervisor_meeting_2026_08_31_v1.pptx)
Prereq: python scripts/deck/make_deck_figures.py   (clean figures in figures/deck_2026_08_31/)
--render DIR renders every slide to PNG via spire.presentation for visual audit (guardrail B1) and
prints the per-slide word / numeric-token / notes report + a banned-string scan.
"""
import argparse
import importlib
import os
import re
import shutil
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
from pptx import Presentation  # noqa: E402
from deck import config as C  # noqa: E402
from deck import helpers as H  # noqa: E402

PARTS = ["deck.slides_answers", "deck.slides_theory", "deck.slides_measure",
         "deck.slides_results", "deck.slides_close", "deck.slides_appendix"]
BANNED = ["0/40", "‖ΔW‖/‖W₀‖", "0.226", "1.07", "ssim_norm 0.6", "confirmed", "settled", "duplication-invari"]
ALLOWED_PHRASES = ["settled on your side"]


def collect(skip_missing):
    fns = []
    for m in PARTS:
        try:
            mod = importlib.import_module(m)
        except ModuleNotFoundError as e:
            if skip_missing and m.split(".")[-1] in str(e):
                print(f"[deck] WARNING: {m} missing — skipped")
                continue
            raise
        fns += list(mod.SLIDES)
    return fns


def audit(pptx, render_dir=None):
    prs = Presentation(pptx)
    bad = []
    for i, s in enumerate(prs.slides, 1):
        txt = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        words = len(txt.split())
        nums = re.findall(r"(?<![A-Za-z_])[-+]?\d+(?:\.\d+)?(?![A-Za-z_])", txt)
        notes = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        low = txt.lower()
        for ph in ALLOWED_PHRASES:
            low = low.replace(ph, "")
        hits = [b for b in BANNED if b.lower() in low]
        flag = "  <-- BANNED " + str(hits) if hits else ""
        if not notes:
            flag += "  <-- NO NOTES"
        print(f"slide {i:2d}: {words:3d} words, {len(nums):2d} numeric tokens, notes {len(notes):4d} chars{flag}")
        if hits or not notes:
            bad.append(i)
    if render_dir:
        render_all(pptx, render_dir)
    return bad


def _subset(src, lo, hi, out):
    """Write a copy of `src` keeping only slides lo..hi (1-based, inclusive) — slide-id-list manipulation (guardrail B9)."""
    prs = Presentation(src)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for i, sid in enumerate(ids, 1):
        if not (lo <= i <= hi):
            prs.part.drop_rel(sid.rId)
            sldIdLst.remove(sid)
    prs.save(out)


def render_all(pptx, render_dir, chunk=10):
    """Render every slide to PNG. spire's free tier converts only the first 10 slides of a file,
    so the deck is rendered in <=10-slide subsets (temporary files next to the renders)."""
    os.makedirs(render_dir, exist_ok=True)
    from spire.presentation import Presentation as SP
    n = len(Presentation(pptx).slides)
    for lo in range(1, n + 1, chunk):
        hi = min(lo + chunk - 1, n)
        tmp = os.path.join(render_dir, f"_chunk_{lo:02d}_{hi:02d}.pptx")
        _subset(pptx, lo, hi, tmp)
        p = SP()
        p.LoadFromFile(tmp)
        for k in range(p.Slides.Count):
            img = p.Slides[k].SaveAsImage()
            img.Save(os.path.join(render_dir, f"slide_{lo + k:02d}.png"))
            img.Dispose()
        p.Dispose()
        os.remove(tmp)
    print(f"[deck] rendered {n} slides -> {render_dir}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--render", default=None)
    ap.add_argument("--skip-missing", action="store_true")
    args = ap.parse_args()
    fns = collect(args.skip_missing)
    prs = Presentation()
    prs.slide_width, prs.slide_height = C.SL_W, C.SL_H
    H.set_total(len(fns))
    for fn in fns:
        fn(prs)
    os.makedirs(os.path.dirname(C.OUT_PPTX), exist_ok=True)
    prs.save(C.OUT_PPTX)
    shutil.copyfile(C.OUT_PPTX, C.OUT_COPY)
    print(f"[deck] {len(fns)} slides -> {C.OUT_PPTX}  (copy: {C.OUT_COPY})")
    bad = audit(C.OUT_PPTX, args.render)
    if bad:
        print(f"[deck] ATTENTION slides {bad}")


if __name__ == "__main__":
    main()
