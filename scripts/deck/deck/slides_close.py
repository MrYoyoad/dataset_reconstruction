"""Part 5 — Close (slide 22 of the 2026-08-31 deck): where we are, and what I need from you."""
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from . import config as C
from . import helpers as H

_Y0 = C.CT + Inches(0.55)


def _world(s, x, y, w, h, letter, name, what, color, fill, placed, placed_color, verdict):
    """One column of the three-worlds strip: header + one-line definition + the experiment placed in it."""
    H.add_rect(s, x, y, w, h, fill_color=fill, line_color=color, line_width=1.25,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pad = Inches(0.14)
    H.add_text(s, f"{letter}  ·  {name}", x + pad, y + Inches(0.1), w - 2 * pad, Inches(0.4), size=15,
               bold=True, font=C.HEAD, color=color)
    H.add_text(s, what, x + pad, y + Inches(0.55), w - 2 * pad, Inches(0.9), size=12, color=C.GRAY)
    # divider
    H.add_line(s, x + pad, y + Inches(1.5), x + w - pad, y + Inches(1.5), color=color, width=0.75)
    H.add_text(s, "experiment placed here", x + pad, y + Inches(1.58), w - 2 * pad, Inches(0.3), size=10,
               italic=True, color=C.LGRAY)
    H.add_text(s, placed, x + pad, y + Inches(1.88), w - 2 * pad, Inches(1.1), size=12, color=C.BLACK)
    H.add_rect(s, x + pad, y + h - Inches(0.55), w - 2 * pad, Inches(0.4), fill_color=C.WHITE,
               line_color=placed_color, line_width=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    H.add_text(s, verdict, x + pad, y + h - Inches(0.55), w - 2 * pad, Inches(0.4), size=12, bold=True,
               color=placed_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _card_small_radius(s, x, y, w, h, title, body, *, color, fill, body_size=13):
    """H.add_card with a small corner radius (tall cards otherwise get a corner that eats the title)."""
    n_before = len(s.shapes)
    H.add_card(s, x, y, w, h, title, body, color=color, fill=fill, body_size=body_size)
    s.shapes[n_before].adjustments[0] = 0.04
    return s.shapes[n_before]


def slide_c1_close(prs):
    s = H.new_slide(prs)
    H.add_title(s, "Where we are, and what I need from you")
    H.add_lead(s, "next: robust adapter-only inversion, with the four controls")
    # ---- left: three-worlds strip -------------------------------------------------------------
    strip_w = Inches(7.9)
    col_gap = Inches(0.18)
    col_w = (strip_w - 2 * col_gap) // 3
    col_h = Inches(3.75)
    y = _Y0 + Inches(0.05)
    H.add_text(s, "three worlds — every negative reconstruction is ambiguous until it names its world",
               C.MX, _Y0 - Inches(0.05), strip_w, Inches(0.3), size=11, italic=True, color=C.GRAY)
    y = _Y0 + Inches(0.3)
    _world(s, C.MX, y, col_w, col_h, "A", "identifiability wall",
           "the information is genuinely not in the adapter", C.RED, C.LIGHTRED_FILL,
           "the ruler: r_J, whitened d², q_eff on col(J)\n→ would establish A when it happens\n(attack-independent, scoped)",
           C.RED, "not our case so far")
    _world(s, C.MX + col_w + col_gap, y, col_w, col_h, "B", "extraction-limited",
           "the information is present; the decoder cannot yet reach it", C.BLUE, C.LIGHTBLUE_FILL,
           "full-gradient ceiling works;\ndirect inversion hits the superposition wall\n→ this is where we are",
           C.BLUE, "we are here")
    _world(s, C.MX + 2 * (col_w + col_gap), y, col_w, col_h, "C", "prior hallucination",
           "the pixels come from the decoder's prior, not the weights", C.AMBER, C.LIGHTGRAY_FILL,
           "disjoint-adapter control:\nsame 'recovery' against an unrelated adapter?\n→ subtract it, always",
           C.AMBER, "excluded (planned control)")
    H.add_text(s, "discipline: A is established by the ruler; C is excluded by the disjoint-adapter control; only what survives both is a real World-B leak",
               C.MX, y + col_h + Inches(0.12), strip_w, Inches(0.55), size=11, color=C.GRAY, italic=True)
    # ---- right: decisions card --------------------------------------------------------------
    rx = C.MX + strip_w + Inches(0.35)
    rw = C.SL_W - C.MX - rx
    _card_small_radius(s, rx, _Y0 + Inches(0.3), rw, Inches(3.3), "decisions",
               "•  SimuDy reframe — agreed on your side?\n\n"
               "•  scale: toy + theory, or ViT / Stable Diffusion?\n\n"
               "•  F5 shared-perturbation compute\n    (honest null at n=8; rotation suggestive)\n\n"
               "•  instance-level atlas zoo — build it?",
               color=C.BLUE, fill=C.LIGHTBLUE_FILL, body_size=13)
    H.add_footer(s)
    H.set_notes(s, """WHAT WE DID: placed every experiment in the three-worlds map (thesis_note_v2 section 4) and listed the four decisions I need.
THE DISCIPLINE IN ONE LINE (thesis_note_v2 section 4): World A is proven by the ruler (attack-independent); World C is excluded by the disjoint-adapter control (prior-independent); only what survives both is a genuine World-B leak. A is a SCOPED guarantee — no local, per-image, linearised information survives under Gaussian seed noise; it is not a guarantee against priors, higher-order effects, or the composition channel (E6 is exactly such an escape).
E7 — THE OPEN MILESTONE (World B): robust adapter-only inversion, minimise ||Y - F(theta_0, x_hat)|| over candidate images. Four controls keep it honest: (i) disjoint-adapter subtraction (leakage above the prior — excludes World C); (ii) staged Jacobians across activations (the activation crux carried through to reconstruction); (iii) local-vs-global init (optimisation failure != non-identifiability); (iv) render the weakest / strongest singular directions as image edits. Full-gradient reconstruction works (SSIM up to ~0.99 MNIST/CIFAR/Flowers; ViT faces 0.38/0.26/0.52); the bridge decoder reaches 0.951 cosine; direct inversion recovers N=4 and superposes at N=10; q_eff shows the directions are present (up to 156/160) — an extraction gap, not a missing-information gap.
DECISIONS: (1) SimuDy reframe (async pointer 2026-06-29) — direct weight inversion overlaps SimuDy; we reframed DI as the known-recipe upper bound; is that settled on your side? (2) Scale: keep to toy + theory (where the ruler is exact) or push to ViT / Stable Diffusion LoRA (where the ruler is only estimable)? (3) F5 shared-perturbation (cross-dataset) compute: honest null at n=8, the rotation variant is suggestive — worth the GPU budget? (4) Instance-level atlas zoo (same digits, different exemplars) — the honest open item after the content-level +0.989; build it?
ANTICIPATED PUSHBACK (thesis_note_v2 section 5, verbatim):
| he'll say | you answer |
| "Detectability isn't reconstruction." | Agreed — the ruler is a noise-free identifiability read, separate from any attack; we quote q_eff and full-gradient recoveries alongside the blurry adapter-only ones. |
| "Your decoder hallucinates a prior." | World C — we subtract it with the disjoint-adapter control and report the difference. |
| "n=24 killed g0." | Honest: +0.857 at n=12 is strong; +0.777 at n=24 is indeterminate (CI [0.53,0.91]). The tercile structure survives; a lead, not a law. |
| "Is this the optimizer, not the map?" | The local-vs-global init control and the permutation null separate optimization failure from genuine non-identifiability. |
| "The composition result is just a t-SNE picture." | No longer just a picture: the cross-fitted cluster-robust test confirms content-level recovery above the recipe baseline (acc-diff +0.989, CI [0.973, 1.005] excludes 0) after fixing a fold bug — the 5 compositions are distinct digit-subsets, so it recovers WHICH DIGITS (content), not the specific instance; graded to ~0 for single-image swaps (arms 0.03-0.07). The gauge confirmation (seed in raw B,A, gone in Delta W) is the second leg. Caveats: G=30 small; CI upper clips >1 (near-ceiling Normal-approx artifact); true instance-level untested. |
PROVENANCE: notes/thesis_note_v2.md sections 3 (E7), 4, 5.""")
    return s


SLIDES = [slide_c1_close]
