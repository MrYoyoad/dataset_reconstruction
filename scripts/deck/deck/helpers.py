"""Slide primitives for the 2026-08-31 deck (python-pptx 1.0.2).

Ported from the May-14 generator (white theme) and extended with:
    set_notes()  — speaker notes (no generator in the repo wrote notes before)
    add_eq()     — place a rendered equation PNG (see eq_render.py)
    add_card()   — bordered card with title + body (native shapes, guardrail B6)
    fit_image()  — place an image inside a box preserving aspect ratio
Every text frame sets word_wrap=True and auto_size=None (guardrail T3).
"""
import os
from PIL import Image
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from . import config as C

_SLIDE_COUNTER = {"n": 0, "total": 0}


def set_total(n):
    _SLIDE_COUNTER["total"] = n


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])      # blank
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = C.WHITE
    _SLIDE_COUNTER["n"] += 1
    return s


def add_text(slide, text, x, y, w, h, *, size=C.SZ_BODY, bold=False, italic=False,
             font=C.BODY, color=C.BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text.split("\n")):
        if i:
            p = tf.add_paragraph()
            p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             size=C.SZ_BODY):
    """Mixed-style text. runs = list of dicts(text, bold, italic, color, size, font) or {'newline': True}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        if r.get("newline"):
            p = tf.add_paragraph()
            p.alignment = align
            continue
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", C.BODY)
        run.font.size = Pt(r.get("size", size))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", C.BLACK)
    return tb


def add_rect(slide, x, y, w, h, *, fill_color=C.WHITE, line_color=C.GRAY, line_width=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.shadow.inherit = False
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width)
    shp.text_frame.text = ""
    return shp


def add_line(slide, x1, y1, x2, y2, *, color=C.GRAY, width=1.0):
    cn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    return cn


def add_arrow(slide, x1, y1, x2, y2, *, color=C.GRAY, width=1.5):
    """Straight connector with an arrow head at the end (OOXML tailEnd)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    cn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return cn


def add_title(slide, text, *, size=C.SZ_TITLE, color=C.BLACK):
    add_text(slide, text, C.MX, C.MY, C.CW, Inches(0.7), size=size, bold=True,
             font=C.HEAD, color=color, anchor=MSO_ANCHOR.MIDDLE)
    rule_y = C.MY + Inches(0.78)
    add_line(slide, C.MX, rule_y, C.SL_W - C.MX, rule_y, color=C.BLUE, width=1.0)


def add_lead(slide, text, *, y=None, size=C.SZ_LEAD, color=C.GRAY, align=PP_ALIGN.LEFT, h=Inches(0.45)):
    """The single line under the title (≤ ~18 words)."""
    if y is None:
        y = C.CT
    return add_text(slide, text, C.MX, y, C.CW, h, size=size, color=color, align=align,
                    anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text=C.FOOTER):
    add_text(slide, text, C.MX, C.SL_H - Inches(0.35), Inches(7.0), Inches(0.25),
             size=C.SZ_FOOT, color=C.LGRAY)
    n, tot = _SLIDE_COUNTER["n"], _SLIDE_COUNTER["total"]
    add_text(slide, f"{n} / {tot}" if tot else f"{n}", C.SL_W - Inches(0.95), C.SL_H - Inches(0.35),
             Inches(0.8), Inches(0.25), size=C.SZ_FOOT, color=C.LGRAY, align=PP_ALIGN.RIGHT)


def add_tag(slide, text, x, y, *, color=C.BLUE, w=Inches(2.6), h=Inches(0.32), size=11):
    """Small rounded label, e.g. the May-ask this slide answers ('your ask: smooth activations')."""
    add_rect(slide, x, y, w, h, fill_color=C.LIGHTBLUE_FILL if color == C.BLUE else C.LIGHTGRAY_FILL,
             line_color=color, line_width=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, text, x + Inches(0.08), y, w - Inches(0.16), h, size=size, bold=True,
             color=color, anchor=MSO_ANCHOR.MIDDLE)


def image_size(path):
    with Image.open(path) as im:
        return im.size


def fit_image(slide, path, x, y, w, h, *, align="center"):
    """Place image inside box (x,y,w,h) preserving aspect; returns (px, py, pw, ph) EMU."""
    iw, ih = image_size(path)
    box_ar = w / h
    img_ar = iw / ih
    if img_ar >= box_ar:
        pw = w
        ph = int(w / img_ar)
    else:
        ph = h
        pw = int(h * img_ar)
    if align == "center":
        px = x + (w - pw) // 2
    elif align == "left":
        px = x
    else:
        px = x + (w - pw)
    py = y + (h - ph) // 2
    slide.shapes.add_picture(path, px, py, pw, ph)
    return px, py, pw, ph


def add_eq(slide, path, x, y, *, w=None, h=None, max_w=None):
    """Place a rendered equation PNG at (x,y). Give ONE of w/h; max_w caps width."""
    iw, ih = image_size(path)
    if h is not None:
        w_ = int(h * iw / ih)
        if max_w is not None and w_ > max_w:
            w_ = max_w
            h = int(w_ * ih / iw)
        return slide.shapes.add_picture(path, x, y, w_, h)
    if w is not None:
        return slide.shapes.add_picture(path, x, y, w, int(w * ih / iw))
    raise ValueError("add_eq needs w or h")


def add_card(slide, x, y, w, h, title, body, *, color=C.BLUE, fill=None, title_size=15, body_size=13,
             body_color=C.BLACK):
    add_rect(slide, x, y, w, h, fill_color=fill or C.WHITE, line_color=color, line_width=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pad = Inches(0.14)
    add_text(slide, title, x + pad, y + Inches(0.08), w - 2 * pad, Inches(0.4), size=title_size,
             bold=True, font=C.HEAD, color=color)
    add_text(slide, body, x + pad, y + Inches(0.5), w - 2 * pad, h - Inches(0.58), size=body_size,
             color=body_color)


def set_notes(slide, text):
    """Speaker notes (plain text). Template: WHAT / WHY THIS FUNCTION / WHY REPRESENTATIVE / GAL-ASK / CAVEATS / PROVENANCE."""
    slide.notes_slide.notes_text_frame.text = text.strip()
